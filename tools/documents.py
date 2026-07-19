import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
from db import NotFound
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

OUTPUT_DIR = "generated"
os.makedirs(OUTPUT_DIR, exist_ok=True)


def generate_invoice_pdf(conn, bill_id, shop_name="Kirana Store", gstin=None):
    bill = conn.execute(
        "SELECT * FROM bills WHERE id=? AND status='finalized'", (bill_id,)
    ).fetchone()
    if not bill:
        raise NotFound(f"finalized bill {bill_id} not found")
    items = conn.execute(
        "SELECT product_name_snapshot, qty, unit_price, gst_rate, cgst_amt, sgst_amt, line_total "
        "FROM bill_items WHERE bill_id=?", (bill_id,),
    ).fetchall()

    path = os.path.join(OUTPUT_DIR, f"invoice_{bill_id}.pdf")
    doc = SimpleDocTemplate(path, pagesize=A4, topMargin=15*mm, bottomMargin=15*mm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"], fontSize=16)
    elements = []

    elements.append(Paragraph(shop_name, title_style))
    if gstin:
        elements.append(Paragraph(f"GSTIN: {gstin}", styles["Normal"]))
    elements.append(Paragraph(f"Tax Invoice — Bill #{bill_id}", styles["Heading2"]))
    elements.append(Paragraph(f"Date: {bill['finalized_at']}", styles["Normal"]))
    if bill["customer_name"]:
        elements.append(Paragraph(f"Customer: {bill['customer_name']}", styles["Normal"]))
    elements.append(Paragraph(f"Payment: {bill['payment_mode'].upper() if bill['payment_mode'] else '-'}"
                               + (f" (ref: {bill['payment_ref']})" if bill["payment_ref"] else ""),
                               styles["Normal"]))
    elements.append(Spacer(1, 8*mm))

    table_data = [["Item", "Qty", "Rate (₹)", "GST%", "CGST (₹)", "SGST (₹)", "Total (₹)"]]
    for it in items:
        table_data.append([
            it["product_name_snapshot"], f"{it['qty']:g}", f"{it['unit_price']:.2f}",
            f"{it['gst_rate']:g}%", f"{it['cgst_amt']:.2f}", f"{it['sgst_amt']:.2f}", f"{it['line_total']:.2f}",
        ])
    table_data.append(["", "", "", "", "Subtotal", "", f"₹{bill['subtotal']:.2f}"])
    table_data.append(["", "", "", "", "CGST", "", f"₹{bill['cgst_total']:.2f}"])
    table_data.append(["", "", "", "", "SGST", "", f"₹{bill['sgst_total']:.2f}"])
    table_data.append(["", "", "", "", "GRAND TOTAL", "", f"₹{bill['grand_total']:.2f}"])

    t = Table(table_data, colWidths=[55*mm, 15*mm, 22*mm, 15*mm, 22*mm, 22*mm, 25*mm])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2d5f3f")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("GRID", (0, 0), (-1, len(items)), 0.5, colors.grey),
        ("FONTNAME", (4, -1), (-1, -1), "Helvetica-Bold"),
        ("LINEABOVE", (4, -4), (-1, -4), 1, colors.black),
        ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
    ]))
    elements.append(t)
    doc.build(elements)
    return {"file_path": path}


def generate_analysis_deck(conn, owner_id, start_date, end_date, shop_name="Kirana Store"):
    """Real charts (matplotlib -> PNG -> embedded), not a screenshot or a
    text dump wearing a .pptx extension."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from tools.reporting import sales_range
    from tools.inventory import low_stock_report

    data = sales_range(conn, owner_id, start_date, end_date)
    stock = low_stock_report(conn, owner_id)

    prs = Presentation()
    blank = prs.slide_layouts[6]

    # --- Slide 1: title ---
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(8.5), Inches(2))
    tf = tb.text_frame
    tf.text = f"{shop_name} — Sales Analysis"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    p2 = tf.add_paragraph()
    p2.text = f"{start_date} to {end_date}"
    p2.font.size = Pt(18)

    # --- Slide 2: KPI summary ---
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tb.text_frame.text = "Summary"
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    kpis = [
        f"Total Sales: ₹{data['total_sales']:.2f}",
        f"Tax Collected: ₹{data['tax_collected']:.2f}",
        f"Bills: {data['bill_count']}",
    ]
    tb2 = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8), Inches(2))
    tf2 = tb2.text_frame
    for i, line in enumerate(kpis):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.text = line
        para.font.size = Pt(20)

    # --- Slide 3: payment mode split chart ---
    if data["by_payment_mode"]:
        fig, ax = plt.subplots(figsize=(6, 4))
        modes = list(data["by_payment_mode"].keys())
        vals = list(data["by_payment_mode"].values())
        ax.pie(vals, labels=modes, autopct="%1.0f%%", colors=["#2d5f3f", "#7bb87f", "#c8e6c9"])
        ax.set_title("Sales by Payment Mode")
        chart_path = os.path.join(OUTPUT_DIR, "chart_payment.png")
        fig.savefig(chart_path, dpi=120, bbox_inches="tight")
        plt.close(fig)

        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(chart_path, Inches(1.5), Inches(0.7), height=Inches(5.5))

    # --- Slide 4: top items chart ---
    if data["top_items"]:
        fig, ax = plt.subplots(figsize=(7, 4))
        names = [i["product"] for i in data["top_items"]]
        qtys = [i["qty_sold"] for i in data["top_items"]]
        ax.barh(names[::-1], qtys[::-1], color="#2d5f3f")
        ax.set_xlabel("Units Sold")
        ax.set_title("Top Selling Items")
        fig.tight_layout()
        chart_path2 = os.path.join(OUTPUT_DIR, "chart_top_items.png")
        fig.savefig(chart_path2, dpi=120, bbox_inches="tight")
        plt.close(fig)

        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(chart_path2, Inches(1), Inches(0.7), height=Inches(5.5))

    # --- Slide 5: stock health ---
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tb.text_frame.text = "Stock Health — Reorder Needed"
    tb.text_frame.paragraphs[0].font.size = Pt(26)
    body = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8), Inches(4))
    tf3 = body.text_frame
    low = stock["low_stock"]
    if not low:
        tf3.text = "All products above reorder level. ✅"
    else:
        for i, item in enumerate(low):
            para = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
            para.text = f"{item['name']}: {item['qty']:g} {item['unit']} left (reorder at {item['reorder_level']:g})"
            para.font.size = Pt(16)

    path = os.path.join(OUTPUT_DIR, f"analysis_{start_date}_{end_date}.pptx")
    prs.save(path)
    return {"file_path": path}
